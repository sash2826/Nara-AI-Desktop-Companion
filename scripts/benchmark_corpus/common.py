"""Shared document-building helpers for the file-organization benchmark corpus.

These helpers wrap reportlab, python-docx, openpyxl, and python-pptx so the
content modules can describe documents declaratively as ordered blocks and let
this module handle the format-specific rendering.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.lib import colors


# Block model -----------------------------------------------------------------
# A document is described as an ordered list of blocks. Each block is a tuple
# ("kind", payload). Supported kinds:
#   ("title", str)
#   ("heading", str)
#   ("para", str)
#   ("bullets", list[str])
#   ("table", (list[str] header, list[list[str]] rows))


Block = tuple


@dataclass
class DocSpec:
    """Declarative description of a single document."""

    filename: str
    blocks: list[Block] = field(default_factory=list)


def _ensure_dir(path: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)


# PDF -------------------------------------------------------------------------
def write_pdf(path: str, blocks: list[Block]) -> None:
    _ensure_dir(path)
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="Body",
            parent=styles["Normal"],
            fontSize=10.5,
            leading=15,
            alignment=TA_LEFT,
            spaceAfter=6,
        )
    )
    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )
    flow = []
    for kind, payload in blocks:
        if kind == "title":
            flow.append(Paragraph(payload, styles["Title"]))
            flow.append(Spacer(1, 0.3 * cm))
        elif kind == "heading":
            flow.append(Spacer(1, 0.2 * cm))
            flow.append(Paragraph(payload, styles["Heading2"]))
        elif kind == "para":
            flow.append(Paragraph(payload, styles["Body"]))
        elif kind == "bullets":
            items = [
                ListItem(Paragraph(text, styles["Body"]), leftIndent=10)
                for text in payload
            ]
            flow.append(ListFlowable(items, bulletType="bullet", start="circle"))
            flow.append(Spacer(1, 0.15 * cm))
        elif kind == "table":
            header, rows = payload
            data = [header] + rows
            table = Table(data, hAlign="LEFT")
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f3a5f")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("FONTSIZE", (0, 0), (-1, -1), 8.5),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef2f7")]),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            flow.append(table)
            flow.append(Spacer(1, 0.2 * cm))
    doc.build(flow)


# DOCX ------------------------------------------------------------------------
def write_docx(path: str, blocks: list[Block]) -> None:
    _ensure_dir(path)
    document = Document()
    normal = document.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    for kind, payload in blocks:
        if kind == "title":
            document.add_heading(payload, level=0)
        elif kind == "heading":
            document.add_heading(payload, level=1)
        elif kind == "para":
            document.add_paragraph(payload)
        elif kind == "bullets":
            for text in payload:
                document.add_paragraph(text, style="List Bullet")
        elif kind == "table":
            header, rows = payload
            table = document.add_table(rows=1, cols=len(header))
            table.style = "Light Grid Accent 1"
            for idx, cell_text in enumerate(header):
                table.rows[0].cells[idx].text = cell_text
            for row in rows:
                cells = table.add_row().cells
                for idx, cell_text in enumerate(row):
                    cells[idx].text = str(cell_text)
    document.save(path)


# XLSX ------------------------------------------------------------------------
@dataclass
class SheetSpec:
    """Declarative description of a single worksheet."""

    title: str
    header: list[str]
    rows: list[list]
    intro: list[str] = field(default_factory=list)


def write_xlsx(path: str, sheets: list[SheetSpec]) -> None:
    _ensure_dir(path)
    wb = Workbook()
    wb.remove(wb.active)
    bold = Font(bold=True)
    header_font = Font(bold=True, color="FFFFFF")
    from openpyxl.styles import PatternFill

    header_fill = PatternFill("solid", fgColor="1F3A5F")
    for sheet in sheets:
        ws = wb.create_sheet(title=sheet.title[:31])
        row_cursor = 1
        for line in sheet.intro:
            ws.cell(row=row_cursor, column=1, value=line).font = bold
            row_cursor += 1
        if sheet.intro:
            row_cursor += 1
        header_row = row_cursor
        for col_idx, name in enumerate(sheet.header, start=1):
            cell = ws.cell(row=header_row, column=col_idx, value=name)
            cell.font = header_font
            cell.fill = header_fill
        for r_offset, row in enumerate(sheet.rows, start=1):
            for col_idx, value in enumerate(row, start=1):
                ws.cell(row=header_row + r_offset, column=col_idx, value=value)
        for col_idx, name in enumerate(sheet.header, start=1):
            max_len = max(
                [len(str(name))] + [len(str(row[col_idx - 1])) for row in sheet.rows if col_idx - 1 < len(row)]
            )
            ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(max_len + 4, 48)
    wb.save(path)


# PPTX ------------------------------------------------------------------------
@dataclass
class SlideSpec:
    """Declarative description of a single slide."""

    title: str
    bullets: list[str] = field(default_factory=list)
    subtitle: str = ""


def write_pptx(path: str, title: str, subtitle: str, slides: list[SlideSpec]) -> None:
    _ensure_dir(path)
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = title
    cover.placeholders[1].text = subtitle

    for slide in slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide.title
        body = s.placeholders[1].text_frame
        body.clear()
        if slide.subtitle:
            body.text = slide.subtitle
        for idx, bullet in enumerate(slide.bullets):
            para = body.paragraphs[0] if (idx == 0 and not slide.subtitle) else body.add_paragraph()
            para.text = bullet
            para.font.size = PptPt(16)
    prs.save(path)
